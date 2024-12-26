from fpdf import FPDF
from pptx import Presentation
import os

async def create_pdf(frames: list, output_path: str):
    """Convert frames to PDF"""
    pdf = FPDF()
    
    for frame in frames:
        pdf.add_page()
        pdf.image(frame, x=10, y=10, w=190)
    
    pdf.output(output_path)
    return output_path

async def create_ppt(frames: list, output_path: str):
    """Convert frames to PowerPoint"""
    prs = Presentation()
    
    for frame in frames:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pic = slide.shapes.add_picture(frame, 0, 0)
        
        aspect_ratio = pic.height / pic.width
        pic.width = prs.slide_width
        pic.height = int(pic.width * aspect_ratio)
    
    prs.save(output_path)
    return output_path